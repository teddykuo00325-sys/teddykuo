"""
Create 凌策公司 presentation for AI 領導人擂台大賽
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MUTED = RGBColor(0xA0, 0xB4, 0xCF)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44)
CARD_BG = RGBColor(0x1E, 0x4D, 0x8C)
DARKER_BG = RGBColor(0x12, 0x2B, 0x4F)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def set_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, x, y, w, h, text, font_size=14, bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT, font_name="Microsoft JhengHei"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Microsoft JhengHei",
             space_before=Pt(4), space_after=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    return p


def add_card(slide, x, y, w, h, title, lines, title_color=ACCENT_BLUE):
    add_rect(slide, x, y, w, h, CARD_BG)
    # Title
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.4),
                 title, font_size=13, bold=True, color=title_color)
    # Content
    cy = y + Inches(0.5)
    for line in lines:
        add_text_box(slide, x + Inches(0.15), cy, w - Inches(0.3), Inches(0.25),
                     line, font_size=10, color=LIGHT_GRAY)
        cy += Inches(0.25)


def slide_title_bar(slide, title):
    """Add a consistent title bar at the top of content slides."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DARKER_BG)
    add_rect(slide, Inches(0), Inches(0.85), Inches(1.2), Inches(0.05), ACCENT_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8.5), Inches(0.6),
                 title, font_size=28, bold=True, color=WHITE)


def add_progress_bar(slide, x, y, w, h, pct, bar_color=GREEN, bg_color=RGBColor(0x0F, 0x1F, 0x3D)):
    add_rect(slide, x, y, w, h, bg_color)
    bar_w = int(w * pct / 100)
    if bar_w > 0:
        add_rect(slide, x, y, bar_w, h, bar_color)


# ════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(s1, DARKER_BG)

# Accent bar at top
add_rect(s1, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

# Company name
add_text_box(s1, Inches(0.8), Inches(1.0), Inches(8.4), Inches(1.0),
             "凌策公司", font_size=48, bold=True, color=WHITE)

# Tagline
add_text_box(s1, Inches(0.8), Inches(1.9), Inches(8.4), Inches(0.7),
             "AI Agent 驅動的軟體公司", font_size=30, bold=False, color=LIGHT_BLUE)

# Subtitle
add_text_box(s1, Inches(0.8), Inches(2.8), Inches(8.4), Inches(0.5),
             "一人監管 + AI Agent 協作的 Palantir 式組織",
             font_size=16, color=MUTED)

# Date + event
add_text_box(s1, Inches(0.8), Inches(3.7), Inches(8.4), Inches(0.4),
             "AI 領導人擂台大賽  |  2026-04-16", font_size=14, color=MUTED)

# Bottom accent bar
add_rect(s1, Inches(0), Inches(5.4), SLIDE_W, Inches(0.06), ACCENT_BLUE)
# Small accent rectangle
add_rect(s1, Inches(0.8), Inches(2.55), Inches(1.5), Inches(0.04), ACCENT_BLUE)


# ════════════════════════════════════════════════════
# SLIDE 2: 公司願景
# ════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, DARK_BLUE)
slide_title_bar(s2, "公司願景")

add_text_box(s2, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.8),
             "以 AI Agent 協作取代傳統人力組織",
             font_size=24, bold=True, color=WHITE)
add_text_box(s2, Inches(0.8), Inches(1.85), Inches(8.4), Inches(0.5),
             "實現 100% AI 驅動的軟體公司", font_size=20, color=LIGHT_BLUE)

# Three vision pillars
pillars = [
    ("AI 原生", "從零開始以 AI Agent\n作為核心生產力"),
    ("一人監管", "人類領導人專注於\n策略決策與品質把關"),
    ("規模化", "10 個 AI Agent 覆蓋\n業務、技術、營運"),
]
for i, (title, desc) in enumerate(pillars):
    cx = Inches(0.8 + i * 3.1)
    cy = Inches(2.8)
    add_rect(s2, cx, cy, Inches(2.7), Inches(2.2), CARD_BG)
    add_text_box(s2, cx + Inches(0.2), cy + Inches(0.2), Inches(2.3), Inches(0.5),
                 title, font_size=20, bold=True, color=GOLD)
    add_text_box(s2, cx + Inches(0.2), cy + Inches(0.8), Inches(2.3), Inches(1.2),
                 desc, font_size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 3: 組織架構
# ════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, DARK_BLUE)
slide_title_bar(s3, "組織架構")

# Human leader at top
add_rect(s3, Inches(3.5), Inches(1.1), Inches(3.0), Inches(0.6), ACCENT_BLUE)
add_text_box(s3, Inches(3.5), Inches(1.15), Inches(3.0), Inches(0.5),
             "人類領導人 (CEO)", font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow down
add_text_box(s3, Inches(4.7), Inches(1.7), Inches(0.6), Inches(0.35),
             "▼", font_size=16, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Orchestrator
add_rect(s3, Inches(3.0), Inches(2.05), Inches(4.0), Inches(0.6), RGBColor(0x1E, 0x40, 0x7A))
shape = s3.shapes[-1]
shape.line.color.rgb = ACCENT_BLUE
shape.line.width = Pt(2)
add_text_box(s3, Inches(3.0), Inches(2.1), Inches(4.0), Inches(0.5),
             "Orchestrator Agent (中央調度)", font_size=14, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Arrow down
add_text_box(s3, Inches(4.7), Inches(2.65), Inches(0.6), Inches(0.35),
             "▼", font_size=16, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Three departments
depts = [
    ("業務開發部", "Sales Agent\nMarketing Agent\nBD Agent", GREEN),
    ("技術研發部", "Architect Agent\nFrontend Agent\nBackend Agent\nQA Agent", ACCENT_BLUE),
    ("營運管理部", "Finance Agent\nHR Agent\nLegal Agent", ORANGE),
]
for i, (name, agents, color) in enumerate(depts):
    dx = Inches(0.5 + i * 3.2)
    dy = Inches(3.1)
    add_rect(s3, dx, dy, Inches(2.9), Inches(0.5), color)
    add_text_box(s3, dx, dy + Inches(0.05), Inches(2.9), Inches(0.4),
                 name, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s3, dx, dy + Inches(0.5), Inches(2.9), Inches(1.5), CARD_BG)
    add_text_box(s3, dx + Inches(0.15), dy + Inches(0.6), Inches(2.6), Inches(1.3),
                 agents, font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 4: AI Agent 團隊
# ════════════════════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, DARK_BLUE)
slide_title_bar(s4, "AI Agent 團隊")

agents_data = [
    ("Orchestrator", "中央任務調度與協作", "Claude Opus"),
    ("Sales Agent", "客戶開發與報價管理", "Claude Sonnet"),
    ("Marketing Agent", "市場分析與內容策略", "Claude Sonnet"),
    ("BD Agent", "商業拓展與合作夥伴", "Claude Sonnet"),
    ("Architect Agent", "系統架構與技術決策", "Claude Opus"),
    ("Frontend Agent", "前端開發與 UI/UX", "Claude Sonnet"),
    ("Backend Agent", "後端開發與 API 設計", "Claude Sonnet"),
    ("QA Agent", "測試自動化與品質保證", "Claude Sonnet"),
    ("Finance Agent", "財務分析與成本控制", "Claude Haiku"),
    ("Legal Agent", "法律合規與合約審查", "Claude Sonnet"),
]

# Table header
header_y = Inches(1.05)
add_rect(s4, Inches(0.5), header_y, Inches(9.0), Inches(0.38), ACCENT_BLUE)
cols = [(Inches(0.6), Inches(2.2), "Agent 名稱"),
        (Inches(2.9), Inches(3.5), "職責"),
        (Inches(6.5), Inches(2.8), "使用模型")]
for cx, cw, ct in cols:
    add_text_box(s4, cx, header_y + Inches(0.04), cw, Inches(0.3),
                 ct, font_size=11, bold=True, color=WHITE)

# Table rows
for i, (name, role, model) in enumerate(agents_data):
    ry = Inches(1.45 + i * 0.4)
    bg_c = CARD_BG if i % 2 == 0 else DARK_BLUE
    add_rect(s4, Inches(0.5), ry, Inches(9.0), Inches(0.38), bg_c)
    add_text_box(s4, Inches(0.6), ry + Inches(0.04), Inches(2.2), Inches(0.3),
                 name, font_size=10, bold=True, color=LIGHT_BLUE)
    add_text_box(s4, Inches(2.9), ry + Inches(0.04), Inches(3.5), Inches(0.3),
                 role, font_size=10, color=LIGHT_GRAY)
    add_text_box(s4, Inches(6.5), ry + Inches(0.04), Inches(2.8), Inches(0.3),
                 model, font_size=10, color=MUTED)


# ════════════════════════════════════════════════════
# SLIDE 5: 技術架構 - 在線系統
# ════════════════════════════════════════════════════
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5, DARK_BLUE)
slide_title_bar(s5, "技術架構 — 在線系統")

online_items = [
    ("Claude API", "核心 LLM 推理引擎\n提供高品質自然語言理解與生成", ACCENT_BLUE),
    ("MCP Protocol", "Model Context Protocol\nAgent 間標準化通訊協定", GREEN),
    ("LiteLLM", "多模型路由與負載均衡\n統一 API 介面管理", ORANGE),
    ("RAG + Qdrant", "向量資料庫檢索增強生成\n企業知識庫即時查詢", LIGHT_BLUE),
    ("Claude Code", "AI 輔助程式碼開發\n自動化軟體工程流程", GOLD),
]

for i, (title, desc, color) in enumerate(online_items):
    row = i // 3
    col = i % 3
    if row == 1 and i == 3:
        col = 0
    if i == 4:
        col = 1
    cx = Inches(0.5 + col * 3.15)
    cy = Inches(1.2 + row * 2.1)
    add_rect(s5, cx, cy, Inches(2.85), Inches(1.8), CARD_BG)
    # Color accent bar on left
    add_rect(s5, cx, cy, Inches(0.06), Inches(1.8), color)
    add_text_box(s5, cx + Inches(0.2), cy + Inches(0.15), Inches(2.5), Inches(0.4),
                 title, font_size=16, bold=True, color=color)
    add_text_box(s5, cx + Inches(0.2), cy + Inches(0.65), Inches(2.5), Inches(1.0),
                 desc, font_size=11, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 6: 技術架構 - 離線系統
# ════════════════════════════════════════════════════
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6, DARK_BLUE)
slide_title_bar(s6, "技術架構 — 離線系統")

add_text_box(s6, Inches(0.8), Inches(1.1), Inches(8.4), Inches(0.4),
             "高安全性離線推理環境 — 適用於機密資料處理", font_size=14, color=MUTED)

offline_items = [
    ("vLLM 推理引擎", "高效能本地推理伺服器\n支援大規模模型部署", RED_ACCENT),
    ("Hermes 3 模型", "開源大語言模型\n離線環境完整運行", ORANGE),
    ("加密磁碟", "全磁碟加密技術\n資料靜態加密保護", GREEN),
    ("Air-gapped 網路隔離", "完全斷網的實體隔離\n防止任何外部數據洩漏", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(offline_items):
    cx = Inches(0.5 + (i % 2) * 4.65)
    cy = Inches(1.7 + (i // 2) * 2.0)
    add_rect(s6, cx, cy, Inches(4.3), Inches(1.7), CARD_BG)
    add_rect(s6, cx, cy, Inches(0.06), Inches(1.7), color)
    add_text_box(s6, cx + Inches(0.25), cy + Inches(0.15), Inches(3.8), Inches(0.4),
                 title, font_size=16, bold=True, color=color)
    add_text_box(s6, cx + Inches(0.25), cy + Inches(0.65), Inches(3.8), Inches(0.9),
                 desc, font_size=12, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 7: 商業模式
# ════════════════════════════════════════════════════
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7, DARK_BLUE)
slide_title_bar(s7, "商業模式")

biz_models = [
    ("AI Agent\n系統整合服務", "為企業客戶打造\n客製化 AI Agent 系統\n\n端到端整合部署\n持續維運與優化", ACCENT_BLUE),
    ("軟體開發\n服務", "AI 驅動的高效\n軟體開發專案\n\n全棧開發能力\n快速交付 MVP", GREEN),
    ("AI 顧問\n服務", "企業 AI 轉型\n策略規劃與諮詢\n\nAI 導入評估\n技術路線圖制定", GOLD),
]

for i, (title, desc, color) in enumerate(biz_models):
    cx = Inches(0.5 + i * 3.15)
    cy = Inches(1.2)
    add_rect(s7, cx, cy, Inches(2.85), Inches(3.8), CARD_BG)
    add_rect(s7, cx, cy, Inches(2.85), Inches(0.06), color)
    add_text_box(s7, cx + Inches(0.2), cy + Inches(0.25), Inches(2.45), Inches(0.9),
                 title, font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(s7, cx + Inches(0.2), cy + Inches(1.3), Inches(2.45), Inches(2.3),
                 desc, font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 8: 客戶認證策略
# ════════════════════════════════════════════════════
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8, DARK_BLUE)
slide_title_bar(s8, "客戶認證策略")

clients = [
    ("addwii 艾迪威", "65%", 65, "AI Agent 系統整合專案\n預計合約金額：NT$ 500,000+", GREEN),
    ("microjet 捷銳", "40%", 40, "軟體開發服務合作\n預計合約金額：NT$ 300,000+", ORANGE),
    ("維明資訊", "30%", 30, "AI 顧問服務\n預計合約金額：NT$ 200,000+", LIGHT_BLUE),
]

for i, (name, pct_str, pct_val, desc, color) in enumerate(clients):
    cy = Inches(1.2 + i * 1.4)
    add_rect(s8, Inches(0.5), cy, Inches(9.0), Inches(1.2), CARD_BG)
    # Client name
    add_text_box(s8, Inches(0.7), cy + Inches(0.1), Inches(2.5), Inches(0.4),
                 name, font_size=16, bold=True, color=WHITE)
    # Percentage large
    add_text_box(s8, Inches(7.5), cy + Inches(0.1), Inches(1.8), Inches(0.5),
                 pct_str, font_size=28, bold=True, color=color, align=PP_ALIGN.RIGHT)
    # Description
    add_text_box(s8, Inches(0.7), cy + Inches(0.55), Inches(5.5), Inches(0.5),
                 desc, font_size=10, color=MUTED)
    # Progress bar
    add_progress_bar(s8, Inches(0.7), cy + Inches(1.0), Inches(8.5), Inches(0.1), pct_val, color)


# ════════════════════════════════════════════════════
# SLIDE 9: 七日作戰計劃
# ════════════════════════════════════════════════════
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9, DARK_BLUE)
slide_title_bar(s9, "七日作戰計劃")

days = [
    ("Day 1", "公司成立", "組織架構確立、Agent 部署"),
    ("Day 2", "系統搭建", "技術架構上線、開發環境就緒"),
    ("Day 3", "業務啟動", "客戶名單整理、提案準備"),
    ("Day 4", "客戶接觸", "首輪客戶拜訪與需求訪談"),
    ("Day 5", "提案製作", "客製化方案與報價完成"),
    ("Day 6", "Demo 準備", "產品演示與系統整合測試"),
    ("Day 7", "成果發表", "最終簡報與成果展示"),
]

# Timeline line
add_rect(s9, Inches(0.9), Inches(2.7), Inches(8.2), Inches(0.04), ACCENT_BLUE)

for i, (day, title, desc) in enumerate(days):
    cx = Inches(0.5 + i * 1.3)
    # Dot on timeline
    dot = s9.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.35), Inches(2.58), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_BLUE if i < 6 else GOLD
    dot.line.fill.background()
    # Day label above
    add_text_box(s9, cx, Inches(1.5), Inches(1.1), Inches(0.35),
                 day, font_size=12, bold=True, color=ACCENT_BLUE if i < 6 else GOLD, align=PP_ALIGN.CENTER)
    add_text_box(s9, cx, Inches(1.85), Inches(1.1), Inches(0.35),
                 title, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Description below timeline
    add_text_box(s9, cx, Inches(3.1), Inches(1.1), Inches(0.8),
                 desc, font_size=9, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 10: Token 成本控制
# ════════════════════════════════════════════════════
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10, DARK_BLUE)
slide_title_bar(s10, "Token 成本控制")

strategies = [
    ("Prompt 快取機制", "重複使用的系統提示詞快取\n減少重複 Token 消耗", "省 40%", GREEN),
    ("模型分級調用", "簡單任務用 Haiku\n複雜任務用 Opus", "省 60%", ACCENT_BLUE),
    ("批次處理優化", "非即時任務批次送出\n享受 Batch API 折扣", "省 50%", ORANGE),
    ("RAG 精準檢索", "向量搜尋精準提取\n減少冗長 Context", "省 30%", LIGHT_BLUE),
    ("離線模型備援", "機密資料用本地 Hermes 3\n零雲端費用", "省 100%", GOLD),
]

for i, (title, desc, saving, color) in enumerate(strategies):
    row = i // 3
    col = i % 3
    if row == 1 and i == 3:
        col = 0
    if i == 4:
        col = 1
    cx = Inches(0.5 + col * 3.15)
    cy = Inches(1.15 + row * 2.15)
    add_rect(s10, cx, cy, Inches(2.85), Inches(1.85), CARD_BG)
    add_rect(s10, cx, cy, Inches(0.06), Inches(1.85), color)
    add_text_box(s10, cx + Inches(0.2), cy + Inches(0.12), Inches(2.5), Inches(0.35),
                 title, font_size=14, bold=True, color=WHITE)
    add_text_box(s10, cx + Inches(0.2), cy + Inches(0.55), Inches(2.5), Inches(0.8),
                 desc, font_size=10, color=MUTED)
    # Saving badge
    add_rect(s10, cx + Inches(1.6), cy + Inches(1.4), Inches(1.1), Inches(0.32), color)
    add_text_box(s10, cx + Inches(1.6), cy + Inches(1.42), Inches(1.1), Inches(0.28),
                 saving, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 11: 法律合規
# ════════════════════════════════════════════════════
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s11, DARK_BLUE)
slide_title_bar(s11, "法律合規")

legal_items = [
    ("商業秘密保護", [
        "離線系統處理機密資料",
        "加密磁碟與 Air-gapped 隔離",
        "嚴格存取權限控管",
    ], RED_ACCENT),
    ("第三方權利保護", [
        "AI 生成內容版權聲明",
        "開源授權合規管理",
        "智慧財產權風險評估",
    ], ACCENT_BLUE),
    ("資料隱私", [
        "符合 GDPR/個資法規範",
        "資料最小化收集原則",
        "定期隱私影響評估",
    ], GREEN),
]

for i, (title, points, color) in enumerate(legal_items):
    cx = Inches(0.5 + i * 3.15)
    cy = Inches(1.2)
    add_rect(s11, cx, cy, Inches(2.85), Inches(3.8), CARD_BG)
    add_rect(s11, cx, cy, Inches(2.85), Inches(0.06), color)
    add_text_box(s11, cx + Inches(0.2), cy + Inches(0.25), Inches(2.45), Inches(0.5),
                 title, font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    for j, pt in enumerate(points):
        add_text_box(s11, cx + Inches(0.25), cy + Inches(1.0 + j * 0.7), Inches(2.35), Inches(0.6),
                     pt, font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 12: 核心優勢
# ════════════════════════════════════════════════════
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s12, DARK_BLUE)
slide_title_bar(s12, "核心優勢")

advantages = [
    ("AI 原生架構", "從第一天起就以 AI Agent\n為核心的組織設計\n非傳統公司的 AI 改造", ACCENT_BLUE),
    ("Palantir 模式", "一人監管 + AI 大軍\n最小人力成本\n最大運營效率", GOLD),
    ("雙系統資安", "在線 + 離線雙軌並行\n兼顧效能與安全\nAir-gapped 最高防護", GREEN),
    ("快速交付", "AI Agent 24/7 運作\n開發速度提升 10 倍\n7 天完成傳統 3 個月工作", ORANGE),
]

for i, (title, desc, color) in enumerate(advantages):
    cx = Inches(0.5 + (i % 2) * 4.65)
    cy = Inches(1.2 + (i // 2) * 2.1)
    add_rect(s12, cx, cy, Inches(4.3), Inches(1.8), CARD_BG)
    add_rect(s12, cx, cy, Inches(4.3), Inches(0.06), color)
    add_text_box(s12, cx + Inches(0.25), cy + Inches(0.2), Inches(3.8), Inches(0.4),
                 title, font_size=18, bold=True, color=color)
    add_text_box(s12, cx + Inches(0.25), cy + Inches(0.7), Inches(3.8), Inches(1.0),
                 desc, font_size=12, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 13: Demo 展示
# ════════════════════════════════════════════════════
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s13, DARK_BLUE)
slide_title_bar(s13, "Demo 展示")

add_text_box(s13, Inches(0.8), Inches(1.15), Inches(8.4), Inches(0.5),
             "AI 協作指揮中心 — Dashboard", font_size=20, bold=True, color=LIGHT_BLUE)

# Mock dashboard area
add_rect(s13, Inches(0.5), Inches(1.8), Inches(9.0), Inches(3.4), CARD_BG)
# Inner panels
panels = [
    (Inches(0.7), Inches(2.0), Inches(2.7), Inches(1.4), "Agent 狀態監控",
     "即時監控 10 個 Agent\n的運行狀態、任務進度\n與資源消耗"),
    (Inches(3.6), Inches(2.0), Inches(2.7), Inches(1.4), "任務看板",
     "Kanban 式任務管理\n自動分派、優先級排序\n與進度追蹤"),
    (Inches(6.5), Inches(2.0), Inches(2.8), Inches(1.4), "對話記錄",
     "Agent 間協作對話\n決策過程透明可追溯\n完整審計軌跡"),
]
for px, py, pw, ph, ptitle, pdesc in panels:
    add_rect(s13, px, py, pw, ph, DARKER_BG)
    add_text_box(s13, px + Inches(0.1), py + Inches(0.08), pw - Inches(0.2), Inches(0.3),
                 ptitle, font_size=12, bold=True, color=ACCENT_BLUE)
    add_text_box(s13, px + Inches(0.1), py + Inches(0.45), pw - Inches(0.2), Inches(0.9),
                 pdesc, font_size=10, color=MUTED)

# Bottom panel
add_rect(s13, Inches(0.7), Inches(3.6), Inches(8.6), Inches(1.3), DARKER_BG)
add_text_box(s13, Inches(0.9), Inches(3.7), Inches(4.0), Inches(0.3),
             "成本分析儀表板", font_size=12, bold=True, color=ACCENT_BLUE)
add_text_box(s13, Inches(0.9), Inches(4.05), Inches(8.2), Inches(0.7),
             "Token 用量追蹤  |  各 Agent 成本分析  |  預算警報系統  |  ROI 即時計算",
             font_size=11, color=MUTED)


# ════════════════════════════════════════════════════
# SLIDE 14: 謝謝
# ════════════════════════════════════════════════════
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s14, DARKER_BG)

add_rect(s14, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(s14, Inches(0.8), Inches(1.2), Inches(8.4), Inches(1.0),
             "謝謝", font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(s14, Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.6),
             "凌策公司 — AI Agent 驅動的軟體公司",
             font_size=20, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_rect(s14, Inches(3.5), Inches(3.0), Inches(3.0), Inches(0.03), ACCENT_BLUE)

add_text_box(s14, Inches(1.5), Inches(3.4), Inches(7.0), Inches(0.4),
             "AI 領導人擂台大賽  |  2026-04-16", font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

add_text_box(s14, Inches(1.5), Inches(4.0), Inches(7.0), Inches(0.4),
             "聯絡方式：lingce.company@email.com", font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

add_rect(s14, Inches(0), Inches(5.4), SLIDE_W, Inches(0.06), ACCENT_BLUE)


# ── Save ──
output_path = r"C:\Users\B00325\Desktop\公司AI比賽\lingce-company\docs\凌策公司_簡報.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
