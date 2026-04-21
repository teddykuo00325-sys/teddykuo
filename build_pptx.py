# -*- coding: utf-8 -*-
"""凌策公司 — 完整功能演示 PPT 產生器
作者：郭祐均 Kuo Yu-Chun
輸出：桌面 凌策公司_完整功能演示_郭祐均_YYYYMMDD_HHMM.pptx
"""
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

AUTHOR_ZH = '郭祐均'
AUTHOR_EN = 'Kuo Yu-Chun'
TS = datetime.now().strftime('%Y%m%d_%H%M')
OUT = fr'C:\Users\B00325\Desktop\凌策公司_完整功能演示_{AUTHOR_ZH}_{TS}.pptx'

# ── 色盤（對齊系統 dashboard 深藍配色）
NAVY        = RGBColor(0x0a, 0x0f, 0x1e)
CARD        = RGBColor(0x0f, 0x17, 0x2a)
SLATE       = RGBColor(0x1e, 0x29, 0x3b)
BORDER      = RGBColor(0x33, 0x41, 0x55)
TEXT        = RGBColor(0xe2, 0xe8, 0xf0)
MUTED       = RGBColor(0x94, 0xa3, 0xb8)
ACCENT_BLUE = RGBColor(0x60, 0xa5, 0xfa)
ACCENT_CYAN = RGBColor(0x22, 0xd3, 0xee)
ACCENT_AMB  = RGBColor(0xfb, 0xbf, 0x24)
ACCENT_EMER = RGBColor(0x34, 0xd3, 0x99)
ACCENT_ROSE = RGBColor(0xfb, 0x71, 0x85)
ACCENT_PURP = RGBColor(0xa7, 0x8b, 0xfa)

FONT = '微軟正黑體'

# ─────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color=NAVY):
    """鋪滿底色"""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.line.fill.background()
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.shadow.inherit = False
    return r


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT, align='left',
             anchor='top', font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE,
                          'bottom': MSO_ANCHOR.BOTTOM}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                   'right': PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_card(slide, x, y, w, h, *, fill=CARD, border=BORDER, border_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.05
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(border_w)
    s.shadow.inherit = False
    return s


def add_pill(slide, x, y, w, h, text, *, fill, color=TEXT, size=11):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    s.text_frame.margin_top = s.text_frame.margin_bottom = Emu(0)
    p = s.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = True
    return s


def page_header(slide, idx, total, section, title):
    """頁頂：頁碼 / 章節 / 大標題"""
    # 頁碼 + section 標
    add_text(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.3),
             f'{section}',
             size=11, color=ACCENT_BLUE)
    add_text(slide, Inches(12.0), Inches(0.3), Inches(1.2), Inches(0.3),
             f'{idx:02d} / {total:02d}',
             size=10, color=MUTED, align='right')
    # 大標題
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12.3), Inches(0.75),
             title, size=32, bold=True, color=TEXT)
    # 藍色分隔線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(1.4),
                                  Inches(0.6), Emu(45720))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()


def footer(slide):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
             f'凌策公司 LingCe  ·  {AUTHOR_ZH} {AUTHOR_EN}  ·  AI 領導人擂台大賽 2026',
             size=9, color=MUTED, align='center')


# ═════════════════════════════════════════════════════════════
TOTAL = 20
slide_idx = 0

# ─── 1. 封面
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
# 漸層感：疊一張半透明藍框
accent_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), SH)
accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = ACCENT_BLUE
accent_bar.line.fill.background()

add_text(s, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
         'LINGCE AI CONSULTING  ·  2026', size=14, color=ACCENT_CYAN)
add_text(s, Inches(1), Inches(2.3), Inches(11), Inches(1.2),
         '凌策公司 LingCe', size=60, bold=True, color=TEXT)
add_text(s, Inches(1), Inches(3.6), Inches(11), Inches(0.6),
         'AI Agent 服務型組織 · 完整功能演示', size=24, color=ACCENT_BLUE)
# 標語
add_text(s, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
         '1 位人類老闆 + 10 位 AI Agent 員工',
         size=18, color=ACCENT_AMB, bold=True)
add_text(s, Inches(1), Inches(5.1), Inches(11), Inches(0.5),
         '用 AI 取代傳統人力，為企業交付 AI 驅動解決方案',
         size=14, color=MUTED)
# 作者資訊
add_text(s, Inches(1), Inches(6.4), Inches(11), Inches(0.4),
         f'作者：{AUTHOR_ZH}（{AUTHOR_EN}）  ·  2026 年 4 月',
         size=12, color=MUTED)


# ─── 2. 作品定位
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '01  DEFINITION', '作品定位 — 什麼是凌策？')

# 中央公式卡
card_y = Inches(1.9)
# 左邊：1 人
lc = add_card(s, Inches(1.2), card_y, Inches(3.3), Inches(3.2),
              fill=RGBColor(0x7c, 0x2d, 0x12), border=ACCENT_AMB)
add_text(s, Inches(1.2), card_y + Inches(0.4), Inches(3.3), Inches(0.6),
         '人類老闆', size=16, color=ACCENT_AMB, align='center', bold=True)
add_text(s, Inches(1.2), card_y + Inches(1.1), Inches(3.3), Inches(1.5),
         '1', size=140, color=TEXT, align='center', bold=True)
add_text(s, Inches(1.2), card_y + Inches(2.6), Inches(3.3), Inches(0.4),
         '唯一實體員工 · 監管/決策/風控', size=11, color=MUTED, align='center')

# 中間：+
add_text(s, Inches(4.6), card_y + Inches(1.1), Inches(0.9), Inches(1.5),
         '+', size=80, color=ACCENT_BLUE, align='center', anchor='middle', bold=True)

# 右邊：10 AI
rc = add_card(s, Inches(5.7), card_y, Inches(3.3), Inches(3.2),
              fill=RGBColor(0x1e, 0x3a, 0x8a), border=ACCENT_BLUE)
add_text(s, Inches(5.7), card_y + Inches(0.4), Inches(3.3), Inches(0.6),
         'AI Agent 員工', size=16, color=ACCENT_BLUE, align='center', bold=True)
add_text(s, Inches(5.7), card_y + Inches(1.1), Inches(3.3), Inches(1.5),
         '10', size=140, color=TEXT, align='center', bold=True)
add_text(s, Inches(5.7), card_y + Inches(2.6), Inches(3.3), Inches(0.4),
         'BD / 客服 / 提案 / 前後端 / QA / 財法務 / 文件', size=11, color=MUTED, align='center')

# = 新世代服務公司
ec = add_card(s, Inches(9.1), card_y, Inches(3.0), Inches(3.2),
              fill=RGBColor(0x4c, 0x1d, 0x95), border=ACCENT_PURP)
add_text(s, Inches(9.1), card_y + Inches(0.4), Inches(3.0), Inches(0.5),
         '=', size=32, color=ACCENT_PURP, align='center', bold=True)
add_text(s, Inches(9.1), card_y + Inches(1.1), Inches(3.0), Inches(0.8),
         '新世代\nIT 服務公司', size=22, color=TEXT, align='center', bold=True)
add_text(s, Inches(9.1), card_y + Inches(2.6), Inches(3.0), Inches(0.4),
         '取代傳統 30+ 人團隊', size=12, color=ACCENT_AMB, align='center')

add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
         '沒有業務、沒有工程師、沒有客服 —— 所有職能皆由 AI Agent 扮演',
         size=16, color=TEXT, align='center')
add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
         '凌策以自身組織親自示範「AI 取代傳統人力」的可行性',
         size=13, color=MUTED, align='center')
footer(s)


# ─── 3. 雙軌商業模式
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '02  BUSINESS MODEL', '雙軌商業模式')

# 軌道一
add_card(s, Inches(0.8), Inches(1.9), Inches(5.8), Inches(4.5),
         fill=RGBColor(0x0c, 0x4a, 0x6e), border=ACCENT_CYAN)
add_text(s, Inches(1.2), Inches(2.1), Inches(5), Inches(0.5),
         '軌道一 · 凌策自用', size=20, color=ACCENT_CYAN, bold=True)
add_text(s, Inches(1.2), Inches(2.7), Inches(5), Inches(0.4),
         'AI Agent 自動接案', size=14, color=MUTED)
for i, line in enumerate([
    '▸ BD Agent 自動評估接案、即時報價',
    '▸ 10 位 AI Agent 並行分工執行',
    '▸ Orchestrator 自動產出進度報告',
    '▸ 每次案件累積知識庫',
    '▸ 證明「1 人 + 10 AI」可服務 140 人規模',
]):
    add_text(s, Inches(1.2), Inches(3.3 + i * 0.45), Inches(5.3), Inches(0.4),
             line, size=13, color=TEXT)

# 軌道二
add_card(s, Inches(6.8), Inches(1.9), Inches(5.8), Inches(4.5),
         fill=RGBColor(0x14, 0x53, 0x2d), border=ACCENT_EMER)
add_text(s, Inches(7.2), Inches(2.1), Inches(5), Inches(0.5),
         '軌道二 · 模組授權', size=20, color=ACCENT_EMER, bold=True)
add_text(s, Inches(7.2), Inches(2.7), Inches(5), Inches(0.4),
         '14 個 AI 模組即買即用', size=14, color=MUTED)
for i, line in enumerate([
    '▸ 月費制 AI 模組商品',
    '▸ 24 小時內凌策協助部署',
    '▸ 持續更新模型與能力',
    '▸ 依客戶需求疊加新模組',
    '▸ 4 項 9 折 / 7 項 85 折',
]):
    add_text(s, Inches(7.2), Inches(3.3 + i * 0.45), Inches(5.3), Inches(0.4),
             line, size=13, color=TEXT)

add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
         '兩軌互補：自用驗證後再授權給客戶',
         size=14, color=ACCENT_BLUE, align='center', bold=True)
footer(s)


# ─── 4. 10 位 AI Agent 員工
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '03  AI AGENTS', '10 位 AI Agent 員工')

agents = [
    ('Orchestrator', '指揮中心', 47, ACCENT_PURP),
    ('BD Agent',       '業務開發', 23, ACCENT_BLUE),
    ('客服 Agent',     '業務開發', 18, ACCENT_BLUE),
    ('提案 Agent',     '業務開發', 12, ACCENT_BLUE),
    ('前端 Agent',     '技術研發', 35, ACCENT_CYAN),
    ('後端 Agent',     '技術研發', 41, ACCENT_CYAN),
    ('QA Agent',       '技術研發', 29, ACCENT_CYAN),
    ('財務 Agent',     '營運管理', 15, ACCENT_EMER),
    ('法務 Agent',     '營運管理', 11, ACCENT_EMER),
    ('文件 Agent',     '營運管理', 20, ACCENT_EMER),
]
for i, (name, dept, tasks, color) in enumerate(agents):
    col = i % 5; row = i // 5
    x = Inches(0.6 + col * 2.48); y = Inches(2.0 + row * 2.3)
    add_card(s, x, y, Inches(2.3), Inches(2.1), fill=CARD, border=color)
    add_text(s, x + Inches(0.1), y + Inches(0.2), Inches(2.1), Inches(0.5),
             name, size=15, bold=True, color=color, align='center')
    add_text(s, x + Inches(0.1), y + Inches(0.8), Inches(2.1), Inches(0.4),
             dept, size=11, color=MUTED, align='center')
    add_text(s, x + Inches(0.1), y + Inches(1.2), Inches(2.1), Inches(0.5),
             str(tasks), size=30, color=TEXT, align='center', bold=True)
    add_text(s, x + Inches(0.1), y + Inches(1.7), Inches(2.1), Inches(0.3),
             '累計任務', size=9, color=MUTED, align='center')

add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.3),
         '全部使用 Ollama Qwen 2.5 7B 本地模型 · 透過 AI 指揮官統一派發',
         size=12, color=ACCENT_AMB, align='center')
footer(s)


# ─── 5. 14 AI 模組商品
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '04  MODULES', '14 個 AI 套裝模組（軌道二授權商品）')

groups = [
    ('通用 (4)',   ACCENT_BLUE, [
        ('AI 文件助理', 'NT$ 3,000'), ('AI 會議記錄', 'NT$ 3,500'),
        ('AI 知識庫',   'NT$ 5,000'), ('AI 客服',     'NT$ 8,000'),
    ]),
    ('製造業 (4)', ACCENT_AMB, [
        ('AI 缺料預警', 'NT$ 6,000'), ('AI 採購建議', 'NT$ 5,500'),
        ('AI 生產排程', 'NT$ 8,000'), ('AI 品質分析', 'NT$ 7,000'),
    ]),
    ('買賣業 (4)', ACCENT_EMER, [
        ('AI 行銷文案', 'NT$ 4,500'), ('AI 客戶分析', 'NT$ 5,500'),
        ('AI 報價助理', 'NT$ 4,000'), ('AI 銷售預測', 'NT$ 6,500'),
    ]),
    ('訂單 (2)',   ACCENT_PURP, [
        ('AI 訂單管理', 'NT$ 5,000'), ('AI 交期預估', 'NT$ 4,000'),
    ]),
]

col_x = [Inches(0.6), Inches(3.7), Inches(6.8), Inches(9.9)]
for ci, (gname, color, items) in enumerate(groups):
    x = col_x[ci]
    add_card(s, x, Inches(1.9), Inches(3.0), Inches(4.6),
             fill=CARD, border=color)
    add_text(s, x, Inches(2.0), Inches(3.0), Inches(0.5),
             gname, size=16, bold=True, color=color, align='center')
    for ii, (name, price) in enumerate(items):
        y = Inches(2.7 + ii * 0.85)
        add_text(s, x + Inches(0.15), y, Inches(2.7), Inches(0.35),
                 name, size=12, color=TEXT, bold=True)
        add_text(s, x + Inches(0.15), y + Inches(0.32), Inches(2.7), Inches(0.32),
                 price + ' / 月', size=11, color=ACCENT_AMB)

add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         '自動折扣：4 項組合 9 折  /  7 項組合 85 折',
         size=13, color=ACCENT_BLUE, align='center', bold=True)
footer(s)


# ─── 6. 三家真實客戶
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '05  CLIENTS', '三家真實服務中客戶')

clients = [
    ('addwii 加我科技', 'B2C 場域無塵室', '6 人組織', '82%',
     ['六款場域產品（嬰兒/廚房/浴室/\n客廳/臥室/餐廳）',
      'HCR-100/200/300 三機型',
      'AI 內容行銷 + 客訴情緒分析',
      'addwiicleanroom.com'],
     ACCENT_BLUE, RGBColor(0x0c, 0x4a, 0x6e)),
    ('MicroJet Technology', 'B2B 精密感測製造', '134 人組織', '88%',
     ['MEMS 壓電微流體 1,600+ 專利',
      'CurieJet P710/P760 感測器',
      'ComeTrue T10/M10 3D 列印',
      '產品 Q&A + CSV 洞察 + PII 合規'],
     ACCENT_AMB, RGBColor(0x78, 0x35, 0x0f)),
    ('維明顧問', '企業顧問業', '評估中', '10%',
     ['商業模式評估中',
      'BD Agent 持續訪談',
      '需求仍在釐清',
      '尚未正式簽約'],
     MUTED, RGBColor(0x33, 0x41, 0x55)),
]
for i, (name, kind, scale, prog, items, color, bg_color) in enumerate(clients):
    x = Inches(0.6 + i * 4.2)
    add_card(s, x, Inches(1.9), Inches(4), Inches(4.8),
             fill=bg_color, border=color)
    add_text(s, x + Inches(0.2), Inches(2.1), Inches(3.6), Inches(0.5),
             kind, size=11, color=MUTED)
    add_text(s, x + Inches(0.2), Inches(2.4), Inches(3.6), Inches(0.6),
             name, size=18, color=TEXT, bold=True)
    add_text(s, x + Inches(0.2), Inches(3.1), Inches(3.6), Inches(0.4),
             f'{scale} · 完成 {prog}', size=11, color=color, bold=True)
    for ii, item in enumerate(items):
        add_text(s, x + Inches(0.2), Inches(3.7 + ii * 0.55), Inches(3.6), Inches(0.55),
                 '▸ ' + item, size=10, color=TEXT)

footer(s)


# ─── 7. ★ 12 坪跨客戶 B2B 整合情境（核心亮點）
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '★ KEY FEATURE',
            '12 坪跨客戶 B2B 整合情境（核心差異化）')

# 左：劇情說明
add_card(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.2),
         fill=RGBColor(0x1c, 0x2b, 0x47), border=ACCENT_AMB)

# 流程：陳先生 → addwii → microjet → 序號綁定 → 安裝
steps = [
    ('陳先生', '終端客戶', '12 坪客廳', ACCENT_ROSE),
    ('addwii', 'B2C 品牌', 'NT$ 168,000', ACCENT_BLUE),
    ('microjet', 'B2B 供貨', 'NT$ 8,800', ACCENT_AMB),
    ('感測器', '自動綁定', '3 組序號', ACCENT_EMER),
    ('到府安裝', '完工交付', '凌策 AI 客服 24x7', ACCENT_PURP),
]
for i, (a, b, c, col) in enumerate(steps):
    x = Inches(0.9 + i * 2.43)
    # 圓形
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.1), Inches(1.8), Inches(1.8))
    circle.fill.solid(); circle.fill.fore_color.rgb = col
    circle.line.fill.background()
    add_text(s, x, Inches(2.5), Inches(1.8), Inches(0.5),
             a, size=14, color=NAVY, align='center', bold=True)
    add_text(s, x, Inches(2.95), Inches(1.8), Inches(0.4),
             b, size=10, color=NAVY, align='center')
    # 下方金額
    add_text(s, x, Inches(4.1), Inches(1.8), Inches(0.4),
             c, size=11, color=col, align='center', bold=True)
    # 箭頭（除最後一個）
    if i < len(steps) - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + Inches(1.85), Inches(2.85),
                                 Inches(0.55), Inches(0.3))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT_BLUE
        arr.line.fill.background()

# 底部說明
add_text(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(0.5),
         '陳先生線上下單 → addwii 向 microjet 下 PO-ADW-2026-0117 → 出貨自動生成 3 組感測器序號',
         size=13, color=TEXT, align='center')
add_text(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.5),
         '→ 綁定至客戶住址（永久追溯）→ 到府安裝 → 凌策 AI 客服 24x7 全程待命',
         size=13, color=TEXT, align='center')

add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6),
         '★ 前端三個視角可即時切換 · 後端狀態機同步稽核 log · 可互動推進里程碑',
         size=14, color=ACCENT_AMB, align='center', bold=True)
footer(s)


# ─── 8. 三視角客戶介面
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '06  PERSPECTIVES', '三視角客戶工作台（模擬客戶端）')

views = [
    ('addwii 工作台', 'B2C 品牌視角', ACCENT_BLUE,
     ['六款場域產品線總覽', '進行中訂單表格', '向 microjet 採購入口',
      '推進里程碑按鈕', 'AI 能力連結']),
    ('microjet 工作台', 'B2B 供貨視角', ACCENT_AMB,
     ['四大產品線展示', '待出貨 B2B 採購單', '「出貨」按鈕自動生成序號',
      '序號綁定紀錄表', '產品型錄']),
    ('終端客戶門戶', '陳先生視角', ACCENT_ROSE,
     ['12 坪客廳案全貌', '8 階段安裝里程碑時間軸', '已綁定感測器序號',
      '產品規格 chip', '下載報價 PDF']),
]
for i, (title, desc, color, items) in enumerate(views):
    x = Inches(0.6 + i * 4.2)
    add_card(s, x, Inches(1.9), Inches(4), Inches(4.9),
             fill=CARD, border=color)
    add_text(s, x + Inches(0.2), Inches(2.1), Inches(3.6), Inches(0.5),
             desc, size=11, color=MUTED)
    add_text(s, x + Inches(0.2), Inches(2.4), Inches(3.6), Inches(0.7),
             title, size=19, color=color, bold=True)
    for ii, it in enumerate(items):
        add_text(s, x + Inches(0.2), Inches(3.3 + ii * 0.55), Inches(3.6), Inches(0.5),
                 '▸ ' + it, size=11, color=TEXT)

add_text(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.3),
         '完整還原 B2B2C 商業關係 · 側邊欄一鍵切換 · 狀態即時同步',
         size=12, color=ACCENT_BLUE, align='center')
footer(s)


# ─── 9. 對外官網
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '07  WEBSITE', '對外官網（潛客入口）')

add_card(s, Inches(0.5), Inches(1.9), Inches(5.8), Inches(4.7),
         fill=CARD, border=ACCENT_CYAN)
add_text(s, Inches(0.8), Inches(2.1), Inches(5.3), Inches(0.5),
         '網站功能（index.html）', size=16, color=ACCENT_CYAN, bold=True)
for i, line in enumerate([
    '▸ Hero 區：1 人 + 10 AI 視覺公式',
    '▸ 雙軌商業模式說明',
    '▸ 14 AI 模組展示 + 即時估價計算器',
    '▸ 三家真實客戶案例',
    '▸ 12 坪跨客戶整合情境',
    '▸ 潛客需求表單',
    '▸ 側邊欄常駐入口「🔐 內部控制台」',
]):
    add_text(s, Inches(0.8), Inches(2.7 + i * 0.48), Inches(5.5), Inches(0.4),
             line, size=12, color=TEXT)

add_card(s, Inches(6.7), Inches(1.9), Inches(6.1), Inches(4.7),
         fill=RGBColor(0x14, 0x53, 0x2d), border=ACCENT_EMER)
add_text(s, Inches(7.0), Inches(2.1), Inches(5.5), Inches(0.5),
         '閉環：潛客 → CRM 詢問單', size=16, color=ACCENT_EMER, bold=True)
add_text(s, Inches(7.0), Inches(2.8), Inches(5.5), Inches(0.4),
         '① 潛客造訪官網', size=13, color=TEXT)
add_text(s, Inches(7.0), Inches(3.2), Inches(5.5), Inches(0.4),
         '② 選擇感興趣模組（即時估價）', size=13, color=TEXT)
add_text(s, Inches(7.0), Inches(3.6), Inches(5.5), Inches(0.4),
         '③ 填寫聯絡資訊 + 需求描述', size=13, color=TEXT)
add_text(s, Inches(7.0), Inches(4.0), Inches(5.5), Inches(0.4),
         '④ 送出 → POST /api/inquiry/submit', size=13, color=TEXT)
add_text(s, Inches(7.0), Inches(4.4), Inches(5.5), Inches(0.4),
         '⑤ 自動寫入 JSONL（不可刪除）', size=13, color=TEXT)
add_text(s, Inches(7.0), Inches(4.8), Inches(5.5), Inches(0.4),
         '⑥ 同步寫入 SQLite CRM 詢問單', size=13, color=TEXT)
add_text(s, Inches(7.0), Inches(5.2), Inches(5.5), Inches(0.4),
         '⑦ 老闆端「📥 潛客詢問」即時看到', size=13, color=TEXT)
add_text(s, Inches(7.0), Inches(5.9), Inches(5.5), Inches(0.4),
         '→ 端對端商業閉環', size=14, color=ACCENT_AMB, bold=True)

footer(s)


# ─── 10. AI 指揮官
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '08  COMMANDER', 'AI 指揮官 — 自然語言派發')

add_text(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.5),
         '老闆用一句話下達指令，Orchestrator 自動分析意圖並派發',
         size=14, color=MUTED, align='center')

# 三層視覺
# 上：老闆說的話
add_card(s, Inches(2), Inches(2.6), Inches(9.3), Inches(0.9),
         fill=RGBColor(0x4c, 0x1d, 0x95), border=ACCENT_PURP)
add_text(s, Inches(2), Inches(2.75), Inches(9.3), Inches(0.6),
         '「幫我分析 addwii 這個月的客訴，找出前 3 名主要問題」',
         size=16, color=TEXT, align='center', bold=True)

# 箭頭
arr1 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), Inches(3.65),
                          Inches(0.7), Inches(0.4))
arr1.fill.solid(); arr1.fill.fore_color.rgb = ACCENT_BLUE
arr1.line.fill.background()

# 中：Orchestrator
add_card(s, Inches(2), Inches(4.15), Inches(9.3), Inches(0.85),
         fill=RGBColor(0x1e, 0x40, 0xaf), border=ACCENT_BLUE)
add_text(s, Inches(2), Inches(4.3), Inches(9.3), Inches(0.5),
         'Orchestrator 意圖分析', size=14, color=ACCENT_CYAN, align='center', bold=True)
add_text(s, Inches(2), Inches(4.7), Inches(9.3), Inches(0.3),
         '→ 辨識為「客戶回饋分析」場景 → 派發至對應 Agent',
         size=11, color=TEXT, align='center')

# 箭頭
arr2 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), Inches(5.15),
                          Inches(0.7), Inches(0.4))
arr2.fill.solid(); arr2.fill.fore_color.rgb = ACCENT_BLUE
arr2.line.fill.background()

# 下：執行結果
add_card(s, Inches(2), Inches(5.65), Inches(9.3), Inches(1.3),
         fill=RGBColor(0x14, 0x53, 0x2d), border=ACCENT_EMER)
add_text(s, Inches(2), Inches(5.8), Inches(9.3), Inches(0.4),
         '自動執行', size=14, color=ACCENT_EMER, align='center', bold=True)
add_text(s, Inches(2), Inches(6.15), Inches(9.3), Inches(0.4),
         '▸ 客服 Agent 讀取客訴紀錄 ▸ 情緒分類 ▸ 產出 Top 3 問題摘要',
         size=11, color=TEXT, align='center')
add_text(s, Inches(2), Inches(6.5), Inches(9.3), Inches(0.4),
         '▸ 全程寫入稽核 log ▸ PII 自動遮蔽',
         size=11, color=TEXT, align='center')

footer(s)


# ─── 11. 客戶驗收中心
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '09  ACCEPTANCE', '客戶驗收中心 — 6 大 AI 能力場景')

scenarios = [
    ('產品 Q&A', 'ChromaDB RAG\n+ bigram fallback', ACCENT_BLUE),
    ('客戶回饋分析', '情緒標註\n+ 自動派工', ACCENT_ROSE),
    ('B2B 提案生成', '規格驗證\n+ Qwen 文案', ACCENT_AMB),
    ('內容行銷', 'FB/IG/Blog\n文案', ACCENT_EMER),
    ('CSV 洞察', '10 房 30 天\n43 萬筆', ACCENT_CYAN),
    ('PII 合規', '9 類偵測\n+ 稽核', ACCENT_PURP),
]
for i, (name, desc, color) in enumerate(scenarios):
    col = i % 3; row = i // 3
    x = Inches(0.7 + col * 4.15); y = Inches(2.0 + row * 2.4)
    add_card(s, x, y, Inches(3.9), Inches(2.2), fill=CARD, border=color)
    add_text(s, x + Inches(0.2), y + Inches(0.3), Inches(3.5), Inches(0.6),
             name, size=18, color=color, bold=True, align='center')
    add_text(s, x + Inches(0.2), y + Inches(1.1), Inches(3.5), Inches(1),
             desc, size=12, color=TEXT, align='center')

add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
         '依客戶 .docx 驗收標準逐項對應 · 每題可實際執行驗收',
         size=12, color=ACCENT_AMB, align='center')
footer(s)


# ─── 12. AI Agent 工作流軌跡
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '10  TRACE', 'AI Agent 工作流軌跡（驗收截圖證據）')

# 三欄
cols = [
    ('① Orchestrator 路由', ACCENT_PURP, [
        '意圖分類：product_qa',
        '派發至：客服 Agent (cs-001)',
        '客戶：addwii',
        '檢索模式：RAG',
    ]),
    ('② RAG 向量檢索 Top-3', ACCENT_BLUE, [
        '8 坪嬰兒房推薦 (0.639)',
        '坪數推薦型號 (0.603)',
        'HCR 系列完整產品線 (0.501)',
        '',
    ]),
    ('③ LLM 呼叫', ACCENT_EMER, [
        'Model: qwen2.5:7b',
        'Elapsed: 17.3 s',
        'PII 遮蔽：0 筆',
        '狀態：成功',
    ]),
]
for i, (title, color, items) in enumerate(cols):
    x = Inches(0.6 + i * 4.2)
    add_card(s, x, Inches(1.9), Inches(4), Inches(4.5),
             fill=CARD, border=color)
    add_text(s, x + Inches(0.2), Inches(2.05), Inches(3.6), Inches(0.55),
             title, size=14, color=color, bold=True)
    for ii, it in enumerate(items):
        if not it: continue
        add_text(s, x + Inches(0.2), Inches(2.8 + ii * 0.55), Inches(3.6), Inches(0.5),
                 it, size=12, color=TEXT)

add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
         '★ 對應驗收規則：「答題須附 AI Agent 工作流截圖」— 不會被當成聊天視窗扣 30%',
         size=12, color=ACCENT_AMB, align='center', bold=True)
footer(s)


# ─── 13. 組織出缺勤
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '11  ATTENDANCE', '組織出缺勤系統（140 人真實 HR）')

add_text(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.5),
         '匯入 microjet 134 人 + addwii 6 人真實組織，含請假/加班/多級審批完整流程',
         size=14, color=MUTED, align='center')

features = [
    ('出缺勤狀態機', '打卡 / 請假 / 加班\n異常預警', ACCENT_BLUE),
    ('多級審批', '經理 → 指定HR\n→ 人事HR → 總經理', ACCENT_AMB),
    ('職務代理人', '請假必須指定\n強制流程', ACCENT_EMER),
    ('統計分析', '日/月報表\nHR 編輯審批', ACCENT_PURP),
    ('階級感知聊天', '跨部門協作\n上署同意機制', ACCENT_CYAN),
    ('編輯稽核', 'append-only\nJSONL 留存', ACCENT_ROSE),
]
for i, (name, desc, color) in enumerate(features):
    col = i % 3; row = i // 3
    x = Inches(0.7 + col * 4.15); y = Inches(2.5 + row * 2.0)
    add_card(s, x, y, Inches(3.9), Inches(1.8), fill=CARD, border=color)
    add_text(s, x + Inches(0.2), y + Inches(0.25), Inches(3.5), Inches(0.5),
             name, size=15, color=color, bold=True, align='center')
    add_text(s, x + Inches(0.2), y + Inches(0.85), Inches(3.5), Inches(0.9),
             desc, size=12, color=TEXT, align='center')

footer(s)


# ─── 14. CRM 客戶管理
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '12  CRM', '客戶管理中心（詢問 → 完工）')

# 五階段 pipeline
stages = [
    ('① 客戶詢問', '新詢問單', ACCENT_BLUE),
    ('② 發報價',   '報價單 PDF', ACCENT_CYAN),
    ('③ 下訂單',   '付訂金', ACCENT_AMB),
    ('④ 安裝中',   '派工執行', ACCENT_PURP),
    ('⑤ 已完成',   '尾款 + 結案', ACCENT_EMER),
]
for i, (name, desc, color) in enumerate(stages):
    x = Inches(0.3 + i * 2.63)
    add_card(s, x, Inches(2.2), Inches(2.3), Inches(1.8),
             fill=CARD, border=color)
    add_text(s, x + Inches(0.1), Inches(2.4), Inches(2.1), Inches(0.55),
             name, size=14, color=color, bold=True, align='center')
    add_text(s, x + Inches(0.1), Inches(3.1), Inches(2.1), Inches(0.6),
             desc, size=11, color=TEXT, align='center')
    if i < 4:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + Inches(2.38), Inches(2.95),
                                 Inches(0.25), Inches(0.3))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT_BLUE
        arr.line.fill.background()

# 下方：6 項 KPI
add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
         '6 項即時 KPI', size=16, color=ACCENT_CYAN, bold=True)
kpis = ['累計客戶詢問', '已發報價單', '成交訂單', '安裝派工中', '本月實收月費', '洽談中案件']
for i, k in enumerate(kpis):
    x = Inches(0.6 + i * 2.06)
    add_card(s, x, Inches(5.1), Inches(1.9), Inches(1.3),
             fill=SLATE, border=BORDER)
    add_text(s, x, Inches(5.3), Inches(1.9), Inches(0.4),
             k, size=10, color=MUTED, align='center')
    add_text(s, x, Inches(5.7), Inches(1.9), Inches(0.6),
             '—', size=22, color=ACCENT_AMB, bold=True, align='center')

add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         '所有狀態改動自動同步關聯單據 · 完整稽核留存',
         size=12, color=ACCENT_AMB, align='center')
footer(s)


# ─── 15. ★ 合規中心（一票否決保命）
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '★ P0 COMPLIANCE',
            '合規中心（驗收構面五 · 一票否決題保命）')

# 四大控制點
controls = [
    ('C1 · 個資不外流',        '雲端 API 已關閉\nCLAUDE_API_DISABLED = True'),
    ('C2 · PII 遮蔽中介層',    'LLM 呼叫前強制\n9 類偵測 + 遮蔽'),
    ('C3 · 稽核日誌完整性',    '4 個 JSONL 檔\nappend-only 不可刪除'),
    ('C4 · 人工審核閘',        '刪除/匯出二次確認\n必填 10 字理由'),
]
for i, (name, desc) in enumerate(controls):
    col = i % 2; row = i // 2
    x = Inches(0.6 + col * 6.3); y = Inches(1.9 + row * 1.7)
    add_card(s, x, y, Inches(6.1), Inches(1.5),
             fill=RGBColor(0x14, 0x53, 0x2d), border=ACCENT_EMER)
    add_text(s, x + Inches(0.25), y + Inches(0.2), Inches(5.7), Inches(0.5),
             '[v]  ' + name, size=15, color=ACCENT_EMER, bold=True)
    add_text(s, x + Inches(0.25), y + Inches(0.75), Inches(5.7), Inches(0.7),
             desc, size=12, color=TEXT)

# PII 9 類
add_card(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.5),
         fill=CARD, border=ACCENT_AMB)
add_text(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5),
         'PII 自動偵測類型（9 類）', size=14, color=ACCENT_AMB, bold=True)
add_text(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.5),
         '身分證字號 · 台灣手機/市話 · Email · 信用卡 · 中文姓名 ·',
         size=11, color=TEXT)
add_text(s, Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.5),
         '英文姓名 · 台灣住址 · roomId · houseId',
         size=11, color=TEXT)

add_text(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
         '即時展示頁：輸入含 PII 文字 → 看到遮蔽前後對比 + 寫入稽核 log',
         size=11, color=ACCENT_CYAN, align='center')
footer(s)


# ─── 16. 三種報價單 PDF
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '13  PDF', '三種報價單模板（正確反映三種商業關係）')

pdfs = [
    ('凌策 AI 授權', 'AI 模組月費', '14 個模組\n4/7 項折扣\n深藍色調', ACCENT_BLUE),
    ('addwii B2C', '場域無塵室整機', 'BOM 清單\n里程碑時程\n訂金尾款\n天藍色調', ACCENT_CYAN),
    ('microjet B2B', '零組件出貨單', '感測器序號綁定\n追溯最終客戶\n琥珀色調', ACCENT_AMB),
]
for i, (title, kind, desc, color) in enumerate(pdfs):
    x = Inches(0.6 + i * 4.2)
    add_card(s, x, Inches(1.9), Inches(4), Inches(4.8),
             fill=CARD, border=color)
    add_text(s, x + Inches(0.2), Inches(2.1), Inches(3.6), Inches(0.4),
             kind, size=11, color=MUTED)
    add_text(s, x + Inches(0.2), Inches(2.5), Inches(3.6), Inches(0.7),
             title, size=20, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(3.4), Inches(3.6), Inches(3),
             desc, size=13, color=TEXT)

add_text(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.3),
         'reportlab + Microsoft 正黑體 · 繁中零依賴 · 評審點按鈕即可下載',
         size=12, color=ACCENT_AMB, align='center')
footer(s)


# ─── 17. 技術堆疊
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '14  TECH STACK', '技術堆疊')

stack = [
    ('前端',         'HTML + Tailwind CSS (CDN)',              ACCENT_CYAN),
    ('後端',         'Flask + Werkzeug (IPv6 dual-stack)',     ACCENT_BLUE),
    ('本地 LLM',     'Ollama + Qwen 2.5 7B  (完全離線)',       ACCENT_PURP),
    ('向量檢索',     'ChromaDB + paraphrase-multilingual',     ACCENT_AMB),
    ('資料儲存',     'SQLite + JSON + JSONL',                  ACCENT_EMER),
    ('PDF 產出',     'reportlab + Microsoft 正黑體',           ACCENT_ROSE),
    ('PII Guard',    '9 種 regex 中介層',                      ACCENT_CYAN),
    ('稽核機制',     'JSONL append-only 非同步寫入',           ACCENT_BLUE),
]
for i, (k, v, color) in enumerate(stack):
    col = i % 2; row = i // 2
    x = Inches(0.6 + col * 6.25); y = Inches(2.0 + row * 1.15)
    add_card(s, x, y, Inches(6.1), Inches(0.95),
             fill=CARD, border=color)
    add_text(s, x + Inches(0.25), y + Inches(0.12), Inches(1.8), Inches(0.35),
             k, size=13, color=color, bold=True)
    add_text(s, x + Inches(0.25), y + Inches(0.5), Inches(5.7), Inches(0.4),
             v, size=12, color=TEXT)

add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
         '全部離線可跑 · 不依賴雲端 API · 符合企業級資料合規',
         size=13, color=ACCENT_AMB, align='center', bold=True)
footer(s)


# ─── 18. 與對手差異化
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '15  DIFFERENTIATION',
            '與其他參賽作品差異化')

compare = [
    ('1 人 + 10 AI 組織架構',         '本作品首創',        '-'),
    ('10 位 AI Agent 並行（非單一）', '本作品',            '-'),
    ('三視角客戶工作台切換',           '本作品',            '-'),
    ('12 坪跨客戶 B2B 整合情境',      '本作品',            '-'),
    ('HCR 坪數自動選型 + CADR 驗證',  '本作品',            '規格留白'),
    ('PII Guard 9 類實作+稽核',       '本作品',            '設計目標'),
    ('三種報價單 PDF 分離',           '本作品',            '單一模板'),
    ('對外官網 + 潛客轉 CRM 閉環',    '本作品',            '-'),
]
# 表頭
add_card(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(0.55),
         fill=RGBColor(0x1e, 0x40, 0xaf), border=ACCENT_BLUE)
add_text(s, Inches(0.8), Inches(1.98), Inches(5), Inches(0.4),
         '創新點', size=13, color=TEXT, bold=True)
add_text(s, Inches(6.5), Inches(1.98), Inches(3), Inches(0.4),
         '本作品（凌策）', size=13, color=ACCENT_EMER, bold=True, align='center')
add_text(s, Inches(9.5), Inches(1.98), Inches(3), Inches(0.4),
         '對手作品（參考）', size=13, color=ACCENT_AMB, bold=True, align='center')

for i, (name, ours, theirs) in enumerate(compare):
    y = Inches(2.55 + i * 0.52)
    bg_color = CARD if i % 2 == 0 else SLATE
    add_card(s, Inches(0.6), y, Inches(12.1), Inches(0.5),
             fill=bg_color, border=BORDER)
    add_text(s, Inches(0.8), y + Inches(0.1), Inches(5.5), Inches(0.35),
             name, size=11, color=TEXT)
    add_text(s, Inches(6.5), y + Inches(0.1), Inches(3), Inches(0.35),
             '[v] ' + ours, size=11, color=ACCENT_EMER, align='center', bold=True)
    add_text(s, Inches(9.5), y + Inches(0.1), Inches(3), Inches(0.35),
             theirs, size=11, color=MUTED, align='center')

footer(s)


# ─── 19. 成果數據
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
page_header(s, slide_idx, TOTAL, '16  METRICS', '成果數據')

kpis = [
    ('10,000+',  '後端核心程式行數', ACCENT_BLUE),
    ('8,500+',   'dashboard 前端行數', ACCENT_CYAN),
    ('10',       'AI Agent 員工', ACCENT_PURP),
    ('6',        '驗收 AI 能力場景', ACCENT_EMER),
    ('14',       'AI 模組商品', ACCENT_AMB),
    ('140',      '真實客戶組織人數', ACCENT_ROSE),
    ('31',       '知識庫 FAQ 筆數', ACCENT_BLUE),
    ('9',        'PII 偵測類型', ACCENT_CYAN),
    ('4',        '稽核 JSONL 檔案', ACCENT_EMER),
    ('3',        '報價單 PDF 模板', ACCENT_AMB),
]
for i, (num, label, color) in enumerate(kpis):
    col = i % 5; row = i // 5
    x = Inches(0.6 + col * 2.48); y = Inches(2.0 + row * 2.2)
    add_card(s, x, y, Inches(2.3), Inches(2.0),
             fill=CARD, border=color)
    add_text(s, x, y + Inches(0.3), Inches(2.3), Inches(1.0),
             num, size=36, color=color, bold=True, align='center')
    add_text(s, x, y + Inches(1.4), Inches(2.3), Inches(0.4),
             label, size=11, color=TEXT, align='center')

footer(s)


# ─── 20. 結語
slide_idx += 1
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), SH)
accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = ACCENT_BLUE
accent_bar.line.fill.background()

add_text(s, Inches(1), Inches(1.6), Inches(11), Inches(0.6),
         'LINGCE · THANK YOU', size=18, color=ACCENT_CYAN)
add_text(s, Inches(1), Inches(2.3), Inches(11), Inches(1.4),
         '以 AI Agent 全面取代人力的', size=36, bold=True, color=TEXT)
add_text(s, Inches(1), Inches(3.1), Inches(11), Inches(1.4),
         '新世代服務公司', size=42, bold=True, color=ACCENT_BLUE)

add_text(s, Inches(1), Inches(4.7), Inches(11), Inches(0.5),
         '凌策不是 demo —— 是真實在運作、穩定服務 140 人規模客戶的系統',
         size=14, color=MUTED)
add_text(s, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
         '每個案件累積的資料、對話、決策，都成為下次模組授權的基礎',
         size=14, color=MUTED)

# 作者區塊
add_card(s, Inches(1), Inches(6.0), Inches(11.3), Inches(1.1),
         fill=CARD, border=ACCENT_BLUE)
add_text(s, Inches(1.3), Inches(6.15), Inches(11), Inches(0.45),
         f'{AUTHOR_ZH} · {AUTHOR_EN}', size=18, color=TEXT, bold=True)
add_text(s, Inches(1.3), Inches(6.65), Inches(11), Inches(0.4),
         f'AI 領導人擂台大賽 · 2026 年 4 月  ·  © 2026 All Rights Reserved',
         size=11, color=MUTED)


# ═════════════════════════════════════════════════════════════
prs.save(OUT)
print(f'Saved: {OUT}')
print(f'Size: {os.path.getsize(OUT) / 1024:.1f} KB')
print(f'Slides: {len(prs.slides)}')
